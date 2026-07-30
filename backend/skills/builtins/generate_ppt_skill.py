"""
PPT 生成 Skill
根据用户提供的主题和大纲，生成 PowerPoint 文件
"""

import json
import logging
import os
import re
import uuid
from typing import Any

from backend.skills.base import BaseSkill

logger = logging.getLogger(__name__)

PPT_OUTPUT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "..", "uploads", "ppt"))

SYSTEM_PROMPT = """你是一位专业的PPT制作专家。你的任务是根据用户提供的主题和要求，生成一份结构清晰、内容专业的PPT大纲和内容。

请按以下JSON格式输出PPT结构：
{
  "title": "PPT标题",
  "subtitle": "副标题",
  "slides": [
    {
      "title": "幻灯片标题",
      "content": ["要点1", "要点2", "要点3"],
      "layout": "title_content"
    }
  ]
}

布局类型(layout)：
- title_slide: 标题页
- title_content: 标题+内容
- two_content: 两栏内容
- section_header: 章节标题页
- comparison: 对比页

要求：
1. 内容专业、简洁，适合演示
2. 每页要点不超过5条
3. 总页数控制在10-20页
4. 输出必须是合法的JSON格式
"""


class GeneratePPTSkill(BaseSkill):
    """生成PPT文件"""

    name = "generate_ppt"
    description = "根据主题生成专业PPT演示文稿。需要主题、页数要求等参数。"
    parameters = {
        "type": "object",
        "properties": {
            "topic": {
                "type": "string",
                "description": "PPT主题",
            },
            "pages": {
                "type": "integer",
                "description": "期望页数（默认10页）",
                "default": 10,
            },
            "audience": {
                "type": "string",
                "description": "目标受众，如'技术团队'、'管理层'、'客户'",
                "default": "通用受众",
            },
            "outline": {
                "type": "string",
                "description": "可选的大纲或要点，用换行分隔",
                "default": "",
            },
        },
        "required": ["topic"],
    }

    async def execute(
        self,
        topic: str,
        pages: int = 10,
        audience: str = "通用受众",
        outline: str = "",
        resume_from: str = "",
        save_style: bool = True,
        **kwargs,
    ) -> str:
        """生成PPT文件（Phase 4.3：风格记忆 + workspace 中间产物 + 可续跑）。"""
        # 兼容 Agent Loop 注入的 user_id / _session_id 等元数据
        try:
            from backend.services.llm import LLMServiceFactory

            style_hint = await self._recall_company_style(
                user_id=kwargs.get("user_id") or kwargs.get("_user_id"),
                identity_id=kwargs.get("_identity_id"),
            )
            # 续跑：从 workspace 中间 JSON 恢复
            ppt_data = None
            if resume_from:
                ppt_data = self._load_intermediate(resume_from)
            if ppt_data is None:
                prompt = (
                    f"请为以下主题生成PPT内容：\n\n主题：{topic}\n"
                    f"目标受众：{audience}\n期望页数：{pages}页\n"
                )
                if outline:
                    prompt += f"大纲要点：\n{outline}\n\n"
                if style_hint:
                    prompt += f"公司风格偏好（必须遵守）：\n{style_hint}\n\n"
                prompt += f"\n请严格按照以下系统提示生成PPT结构：\n{SYSTEM_PROMPT}"

                llm_service = LLMServiceFactory.get_service()
                response = ""
                async for chunk in llm_service.chat(
                    [{"role": "user", "content": prompt}], stream=False
                ):
                    response += chunk.delta or ""
                    if chunk.finish_reason:
                        break

                ppt_data = self._extract_json(response)
                if not ppt_data:
                    return (
                        f"[Error] 无法解析PPT结构。LLM响应：{response[:500]}。"
                        f"下一步：简化主题后重试，或提供更明确的 outline。"
                    )

            # 中间产物落 workspace，失败可 resume_from
            inter_path = self._save_intermediate(ppt_data, topic=topic)

            if save_style and style_hint is None and outline:
                await self._remember_style_hint(
                    f"PPT 风格线索（来自任务）：{outline[:400]}",
                    user_id=kwargs.get("user_id") or kwargs.get("_user_id"),
                    identity_id=kwargs.get("_identity_id"),
                    run_id=kwargs.get("_run_id") or kwargs.get("_agent_run_id"),
                )

            file_path = self._generate_ppt_file(ppt_data)
            file_name = os.path.basename(file_path)
            is_pptx = str(file_path).lower().endswith(".pptx")
            fmt = "PowerPoint .pptx" if is_pptx else "Markdown fallback"
            warn = ""
            if not is_pptx:
                warn = (
                    "\nWARN: not pptx. pip install python-pptx && restart backend"
                )
            parts = [
                "[Success] PPT generated!",
                f"title={ppt_data.get('title', topic)}",
                f"slides={len(ppt_data.get('slides', []))}",
                f"format={fmt}",
                f"path={file_path}",
                f"download=/uploads/ppt/{file_name}",
                f"intermediate={inter_path}",
                "hint=失败可用 resume_from=中间 JSON 路径续跑",
            ]
            if style_hint:
                parts.append("style=applied_from_memory")
            if warn:
                parts.append(warn.strip())
            return "\n".join(parts)
        except Exception as e:
            logger.error(f"PPT generation failed: {e}")
            return (
                f"[Error] PPT生成失败: {type(e).__name__}。"
                f"下一步：检查 python-pptx / 中间 JSON 是否可 resume_from。"
            )

    async def _recall_company_style(
        self, *, user_id=None, identity_id=None
    ) -> str | None:
        try:
            from backend.services import memory_bus

            hits = await memory_bus.recall(
                "PPT 风格 封面 模板 公司",
                kinds=["preference", "methodology", "graph"],
                top_k=3,
                identity_id=identity_id,
                user_id=user_id,
            )
            if not hits:
                return None
            return "\n".join(
                f"- {h.title or h.kind}: {(h.content or '')[:200]}" for h in hits
            )
        except Exception as e:
            logger.debug("ppt style recall skip: %s", e)
            return None

    async def _remember_style_hint(
        self, content: str, *, user_id=None, identity_id=None, run_id=None
    ) -> None:
        try:
            from backend.services import memory_bus

            await memory_bus.remember(
                "preference" if identity_id else "preference",
                content,
                title="PPT 公司风格",
                identity_id=identity_id,
                user_id=user_id,
                source_run_id=run_id,
                source="agent",
                tags=["ppt", "style"],
            )
        except Exception as e:
            logger.debug("ppt style remember skip: %s", e)

    def _workspace_ppt_dir(self) -> str:
        try:
            from backend.tools.permissions import resolve_agent_workspace_root

            root = resolve_agent_workspace_root()
        except Exception:
            root = os.getcwd()
        d = os.path.join(root, ".takton", "ppt_work")
        os.makedirs(d, exist_ok=True)
        return d

    def _save_intermediate(self, ppt_data: dict[str, Any], *, topic: str) -> str:
        d = self._workspace_ppt_dir()
        safe = re.sub(r"[^\w\-]+", "_", (topic or "deck")[:40], flags=re.U)
        path = os.path.join(d, f"{safe}_{uuid.uuid4().hex[:8]}.outline.json")
        with open(path, "w", encoding="utf-8") as f:
            json.dump(ppt_data, f, ensure_ascii=False, indent=2)
        return path

    def _load_intermediate(self, path: str) -> dict[str, Any] | None:
        try:
            p = os.path.abspath(path)
            if not os.path.isfile(p):
                return None
            with open(p, encoding="utf-8") as f:
                data = json.load(f)
            return data if isinstance(data, dict) else None
        except Exception:
            return None

    def _extract_json(self, text: str) -> dict[str, Any] | None:
        """从LLM响应中提取JSON"""
        # 尝试找 ```json 代码块
        import re
        code_block = re.search(r"```(?:json)?\s*(.*?)\s*```", text, re.DOTALL)
        if code_block:
            text = code_block.group(1)

        # 尝试找最外层的大括号
        start = text.find("{")
        end = text.rfind("}")
        if start != -1 and end != -1 and end > start:
            try:
                return json.loads(text[start:end+1])
            except json.JSONDecodeError:
                pass
        return None

    def _generate_ppt_file(self, ppt_data: dict[str, Any]) -> str:
        """使用python-pptx生成PPT文件"""
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError:
            # 如果没有python-pptx，生成markdown文件作为fallback
            return self._generate_markdown_fallback(ppt_data)

        prs = Presentation()
        title = ppt_data.get("title", "Untitled")
        subtitle = ppt_data.get("subtitle", "")
        slides = ppt_data.get("slides", [])

        # 如果没有slides，创建一个默认的
        if not slides:
            slides = [{"title": title, "content": ["内容待补充"], "layout": "title_content"}]

        for i, slide_data in enumerate(slides):
            layout_name = slide_data.get("layout", "title_content")
            slide_title = slide_data.get("title", "")
            content = slide_data.get("content", [])

            # 选择布局
            if layout_name == "title_slide" or i == 0:
                slide_layout = prs.slide_layouts[0]  # Title Slide
            elif layout_name == "section_header":
                slide_layout = prs.slide_layouts[2]  # Section Header
            elif layout_name == "two_content":
                slide_layout = prs.slide_layouts[5]  # Two Content
            else:
                slide_layout = prs.slide_layouts[1]  # Title and Content

            slide = prs.slides.add_slide(slide_layout)

            # 设置标题
            if slide.shapes.title:
                slide.shapes.title.text = slide_title

            # 设置内容
            if layout_name == "title_slide" and i == 0:
                # 标题页设置副标题
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:
                        shape.text = subtitle or ""
            elif layout_name == "two_content" and content and len(slide.placeholders) > 2:
                # 双栏布局：将内容平均拆分到左右两栏
                mid = (len(content) + 1) // 2
                left_items, right_items = content[:mid], content[mid:]
                for idx, items in ((1, left_items), (2, right_items)):
                    body_shape = slide.placeholders[idx]
                    tf = body_shape.text_frame
                    tf.clear()
                    for j, item in enumerate(items):
                        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                        p.text = f"• {item}"
                        p.level = 0
                        p.font.size = Pt(18)
            elif content and len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                for j, item in enumerate(content):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {item}"
                    p.level = 0
                    p.font.size = Pt(18)

        # 保存文件
        os.makedirs(PPT_OUTPUT_DIR, exist_ok=True)
        file_name = f"{uuid.uuid4().hex}.pptx"
        file_path = os.path.join(PPT_OUTPUT_DIR, file_name)
        prs.save(file_path)
        return file_path

    def _generate_markdown_fallback(self, ppt_data: dict[str, Any]) -> str:
        """没有 python-pptx 时的 fallback：生成 markdown，并明确提示安装依赖。"""
        title = ppt_data.get("title", "Untitled")
        subtitle = ppt_data.get("subtitle", "")
        slides = ppt_data.get("slides", [])

        lines = [
            f"# {title}",
            f"\n> {subtitle}\n",
            "\n> ⚠️ 本机未安装 python-pptx，已降级为 Markdown 大纲。\n"
            "> 安装：`pip install python-pptx` 后重试 generate_ppt 可得到 .pptx\n",
        ]
        for slide in slides:
            lines.append(f"\n## {slide.get('title', '')}")
            for item in slide.get("content", []):
                lines.append(f"- {item}")

        os.makedirs(PPT_OUTPUT_DIR, exist_ok=True)
        file_name = f"{uuid.uuid4().hex}.md"
        file_path = os.path.join(PPT_OUTPUT_DIR, file_name)
        with open(file_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))
        return file_path
